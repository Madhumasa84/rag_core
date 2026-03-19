from pypdf import PdfReader
from pptx import Presentation
from bs4 import BeautifulSoup
import os

def extract_from_pdf(path: str) -> str:
    """Extract text from PDF file"""
    reader = PdfReader(path)
    text = ""
    for page in reader.pages:
        extracted = page.extract_text()
        if extracted:
            text += extracted + "\n"
    return text

def extract_from_pptx(path: str) -> str:
    """Extract text from PowerPoint file"""
    try:
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error extracting from PPTX: {e}")
        return ""

def extract_from_html(path: str) -> str:
    """Extract text from HTML file"""
    try:
        with open(path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            return soup.get_text()
    except Exception as e:
        print(f"Error extracting from HTML: {e}")
        return ""

def extract_from_txt(path: str) -> str:
    """Extract text from TXT file"""
    try:
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"Error extracting from TXT: {e}")
        return ""

def extract_text(path: str) -> str:
    """
    Universal extractor - detects file type and extracts text
    Supports: PDF, PPTX, HTML, TXT
    """
    ext = os.path.splitext(path)[1].lower()
    
    if ext == '.pdf':
        return extract_from_pdf(path)
    elif ext in ['.pptx', '.ppt']:
        return extract_from_pptx(path)
    elif ext in ['.html', '.htm']:
        return extract_from_html(path)
    elif ext == '.txt':
        return extract_from_txt(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported: .pdf, .pptx, .ppt, .html, .htm, .txt")